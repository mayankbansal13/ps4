import os
import hashlib
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import fitz  # PyMuPDF
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import pandas as pd
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from config import Config

# Setup logging
logging.basicConfig(level=getattr(logging, Config.LOG_LEVEL))
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles document ingestion, text extraction, and chunking"""
    
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.CHUNK_SIZE,
            chunk_overlap=Config.CHUNK_OVERLAP,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF with OCR fallback for scanned documents"""
        chunks = []
        try:
            # Try PyMuPDF first (for text-based PDFs)
            doc = fitz.open(file_path)
            text_content = ""
            
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                # If page has little to no text, use OCR
                if len(page_text.strip()) < 50:
                    logger.info(f"Using OCR for page {page_num + 1} of {file_path}")
                    page_text = self._ocr_page(file_path, page_num)
                
                if page_text.strip():
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
            
            doc.close()
            
            if text_content.strip():
                text_chunks = self.text_splitter.split_text(text_content)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "pdf"
                        }
                    })
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            # Fallback to full OCR
            chunks = self._ocr_entire_pdf(file_path)
        
        return chunks
    
    def _ocr_page(self, pdf_path: str, page_num: int) -> str:
        """Perform OCR on a specific PDF page"""
        try:
            images = convert_from_path(pdf_path, first_page=page_num+1, last_page=page_num+1)
            if images:
                text = pytesseract.image_to_string(images[0], config=Config.TESSERACT_CONFIG)
                return text
        except Exception as e:
            logger.error(f"OCR failed for page {page_num} of {pdf_path}: {str(e)}")
        return ""
    
    def _ocr_entire_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """Perform OCR on entire PDF as fallback"""
        chunks = []
        try:
            logger.info(f"Performing full OCR on {pdf_path}")
            images = convert_from_path(pdf_path)
            
            full_text = ""
            for i, image in enumerate(images):
                page_text = pytesseract.image_to_string(image, config=Config.TESSERACT_CONFIG)
                if page_text.strip():
                    full_text += f"\n--- Page {i + 1} ---\n{page_text}\n"
            
            if full_text.strip():
                text_chunks = self.text_splitter.split_text(full_text)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(pdf_path),
                            "file_path": pdf_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "pdf",
                            "ocr_processed": True
                        }
                    })
        except Exception as e:
            logger.error(f"Full OCR failed for {pdf_path}: {str(e)}")
        
        return chunks
    
    def extract_text_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from DOCX file"""
        chunks = []
        try:
            doc = Document(file_path)
            text_content = ""
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content += para.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_content += row_text + "\n"
            
            if text_content.strip():
                text_chunks = self.text_splitter.split_text(text_content)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "docx"
                        }
                    })
                    
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
        
        return chunks
    
    def extract_text_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PPTX file"""
        chunks = []
        try:
            prs = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    text_content += slide_text
            
            if text_content.strip():
                text_chunks = self.text_splitter.split_text(text_content)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "pptx"
                        }
                    })
                    
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
        
        return chunks
    
    def extract_text_from_txt(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from TXT file"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text_content = file.read()
            
            if text_content.strip():
                text_chunks = self.text_splitter.split_text(text_content)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "txt"
                        }
                    })
                    
        except Exception as e:
            logger.error(f"Error processing TXT {file_path}: {str(e)}")
        
        return chunks
    
    def extract_text_from_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from CSV file"""
        chunks = []
        try:
            df = pd.read_csv(file_path)
            
            # Convert DataFrame to text representation
            text_content = f"CSV Data from {os.path.basename(file_path)}:\n\n"
            text_content += df.to_string(index=False)
            
            # Also include column descriptions
            text_content += f"\n\nColumns: {', '.join(df.columns.tolist())}\n"
            text_content += f"Total rows: {len(df)}\n"
            
            if text_content.strip():
                text_chunks = self.text_splitter.split_text(text_content)
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "file_path": file_path,
                            "chunk_id": i,
                            "content_hash": self._generate_hash(chunk),
                            "file_type": "csv",
                            "rows": len(df),
                            "columns": df.columns.tolist()
                        }
                    })
                    
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {str(e)}")
        
        return chunks
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Main method to process any supported document type"""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in Config.SUPPORTED_EXTENSIONS:
            logger.error(f"Unsupported file type: {file_ext}")
            return []
        
        logger.info(f"Processing document: {file_path}")
        
        if file_ext == ".pdf":
            return self.extract_text_from_pdf(file_path)
        elif file_ext == ".docx":
            return self.extract_text_from_docx(file_path)
        elif file_ext == ".pptx":
            return self.extract_text_from_pptx(file_path)
        elif file_ext == ".txt":
            return self.extract_text_from_txt(file_path)
        elif file_ext == ".csv":
            return self.extract_text_from_csv(file_path)
        else:
            logger.error(f"Handler not implemented for {file_ext}")
            return []
    
    def process_folder(self, folder_path: str) -> List[Dict[str, Any]]:
        """Process all supported documents in a folder"""
        all_chunks = []
        folder_path = Path(folder_path)
        
        if not folder_path.exists():
            logger.error(f"Folder does not exist: {folder_path}")
            return []
        
        for file_path in folder_path.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in Config.SUPPORTED_EXTENSIONS:
                chunks = self.process_document(str(file_path))
                all_chunks.extend(chunks)
        
        logger.info(f"Processed {len(all_chunks)} chunks from folder: {folder_path}")
        return all_chunks
    
    def _generate_hash(self, content: str) -> str:
        """Generate MD5 hash for content"""
        return hashlib.md5(content.encode()).hexdigest()